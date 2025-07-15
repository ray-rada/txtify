"""
This module provides functionality to convert PowerPoint, Word, and PDF files to plain text.
"""

import os
from io import StringIO
from typing import List, Union

import docx2txt
import pandas as pd
from pptx import Presentation
from PyPDF2 import PdfReader


def convert_pptx_to_text(file_path: str) -> str:
    """
    Convert a PowerPoint file to plain text.

    Args:
        file_path (str): Path to the PowerPoint file.

    Returns:
        str: Extracted text from the PowerPoint file.
    """
    prs = Presentation(file_path)
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    all_text.append(paragraph.text)
    return "\n".join(all_text)


def convert_docx_to_text(file_path: str) -> str:
    """
    Convert a Word document to plain text.

    Args:
        file_path (str): Path to the Word document.

    Returns:
        str: Extracted text from the Word document.
    """
    return "\n".join(docx2txt.process(file_path).splitlines())


def convert_pdf_to_text(file_path: str) -> str:
    """
    Convert a PDF file to plain text.

    Args:
        file_path (str): Path to the PDF file.

    Returns:
        str: Extracted text from the PDF file.
    """
    reader = PdfReader(file_path)
    text = []
    for page in reader.pages:
        text.append(page.extract_text())
    return "\n".join(text)


def convert_xlsx_to_text(file_path: str) -> str:
    buffer = StringIO()
    df = pd.read_excel(file_path)
    df.to_csv(buffer, index=False)
    buffer.seek(0)
    return buffer.read()


def convert_to_text(file_path: str) -> str:
    """
    Convert a file to plain text based on its extension.

    Args:
        file_path (str): Path to the file.

    Returns:
        str: Extracted text from the file.

    Raises:
        ValueError: If the file format is not supported.
    """
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    if ext == ".pptx":
        return convert_pptx_to_text(file_path)
    elif ext == ".docx":
        return convert_docx_to_text(file_path)
    elif ext == ".pdf":
        return convert_pdf_to_text(file_path)
    elif ext == ".xlsx":
        return convert_xlsx_to_text(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")


def get_valid_files(path: str) -> List[str]:
    """
    Get a list of valid files to convert from a given path.

    Args:
        path (str): Path to a file or directory.

    Returns:
        List[str]: List of valid file paths.
    """
    valid_extensions = (".pptx", ".docx", ".pdf", ".xlsx")
    valid_files = []

    if os.path.isfile(path):
        if path.lower().endswith(valid_extensions):
            valid_files.append(path)
    elif os.path.isdir(path):
        for root, _, files in os.walk(path):
            for file in files:
                if file.lower().endswith(valid_extensions):
                    valid_files.append(os.path.join(root, file))

    return valid_files


def batch_convert(input_paths: Union[str, List[str]], output_dir: str) -> None:
    """
    Convert multiple files or directories to plain text and save them in the output directory.

    Args:
        input_paths (Union[str, List[str]]): Path(s) to file(s) or directory(ies) to convert.
        output_dir (str): Directory to save the converted text files.

    Raises:
        OSError: If there's an error creating the output directory.
    """
    os.makedirs(output_dir, exist_ok=True)

    if isinstance(input_paths, str):
        input_paths = [input_paths]

    for input_path in input_paths:
        valid_files = get_valid_files(input_path)

        for file_path in valid_files:
            try:
                text = convert_to_text(file_path)
                rel_path = os.path.relpath(file_path, start=input_path)
                output_file = os.path.join(
                    output_dir, f"{os.path.splitext(rel_path)[0]}.txt"
                )

                os.makedirs(os.path.dirname(output_file), exist_ok=True)

                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(text)

                print(f"Successfully converted {file_path} to {output_file}")
            except Exception as e:
                print(f"Error converting {file_path}: {str(e)}")
